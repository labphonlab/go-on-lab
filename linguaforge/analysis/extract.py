"""Input-format normalization — turns whatever the teacher handed over
(.md/.txt/.html/.rtf/.docx/.pdf/.pptx) into the same markdown-ish plain text
that the rest of the analysis layer already knows how to read (parser.py's
_first_heading / classify.py's section body). One format in, one shape out.

Heavy per-format libraries are imported lazily inside each function so a
teacher who only ever sends Markdown never needs pypdf/python-docx/etc.
installed.
"""

from __future__ import annotations

from html.parser import HTMLParser
from pathlib import Path

SUPPORTED_EXTENSIONS = (".md", ".markdown", ".txt", ".html", ".htm", ".rtf", ".docx", ".pdf", ".pptx")

_HEADING_TAGS = {f"h{n}": "#" * n for n in range(1, 7)}
_BLOCK_TAGS = {"p", "div", "li", "br", "tr", *_HEADING_TAGS}


class _HTMLToMarkdown(HTMLParser):
    """Minimal, dependency-free HTML->text: keeps headings as '#'-prefixed
    lines and list items as '- '-prefixed lines, drops everything else
    (styles, scripts, attributes)."""

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.lines: list[str] = [""]
        self._skip_depth = 0
        self._prefix_stack: list[str] = []

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip_depth += 1
            return
        if tag in _HEADING_TAGS:
            self._flush_line()
            self._prefix_stack.append(_HEADING_TAGS[tag] + " ")
        elif tag == "li":
            self._flush_line()
            self._prefix_stack.append("- ")
        elif tag in _BLOCK_TAGS:
            self._flush_line()

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip_depth = max(0, self._skip_depth - 1)
            return
        if tag in _HEADING_TAGS or tag == "li":
            self._flush_line()
            if self._prefix_stack:
                self._prefix_stack.pop()
        elif tag in _BLOCK_TAGS:
            self._flush_line()

    def handle_data(self, data):
        if self._skip_depth:
            return
        text = " ".join(data.split())
        if text:
            self.lines[-1] += text + " "

    def _flush_line(self):
        if self.lines[-1].strip():
            prefix = self._prefix_stack[-1] if self._prefix_stack else ""
            self.lines[-1] = prefix + self.lines[-1].strip()
            self.lines.append("")

    def result(self) -> str:
        self._flush_line()
        return "\n".join(line for line in self.lines if line.strip())


def _extract_html(path: Path) -> str:
    parser = _HTMLToMarkdown()
    parser.feed(path.read_text(encoding="utf-8", errors="replace"))
    return parser.result()


def _extract_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError as e:
        raise RuntimeError("striprtf is required for .rtf input (pip install striprtf)") from e
    return rtf_to_text(path.read_text(encoding="utf-8", errors="replace"))


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError as e:
        raise RuntimeError("python-docx is required for .docx input (pip install python-docx)") from e

    document = docx.Document(str(path))
    lines = []
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if "heading" in style or "title" in style:
            level = "".join(ch for ch in style if ch.isdigit()) or "1"
            lines.append(f"{'#' * min(int(level), 6)} {text}")
        elif "list" in style or "bullet" in style:
            lines.append(f"- {text}")
        else:
            lines.append(text)
    return "\n".join(lines)


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise RuntimeError("pypdf is required for .pdf input (pip install pypdf)") from e

    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    # PDFs have no structural heading info once flattened to text — the
    # classify step's Claude call (or heuristic fallback) has to infer
    # structure from plain paragraphs alone.
    return "\n\n".join(p.strip() for p in pages if p.strip())


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("python-pptx is required for .pptx input (pip install python-pptx)") from e

    presentation = Presentation(str(path))
    lines = []
    for slide in presentation.slides:
        shapes_text = []
        title_text = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                title_text = text
            else:
                shapes_text.append(text)
        if title_text:
            lines.append(f"# {title_text}")
        for text in shapes_text:
            for line in text.splitlines():
                line = line.strip()
                if line:
                    lines.append(f"- {line}")
    return "\n".join(lines)


_EXTRACTORS = {
    ".html": _extract_html,
    ".htm": _extract_html,
    ".rtf": _extract_rtf,
    ".docx": _extract_docx,
    ".pdf": _extract_pdf,
    ".pptx": _extract_pptx,
}


def extract_text(path: Path) -> str:
    """Normalize any supported input file to markdown-ish plain text."""
    suffix = path.suffix.lower()
    if suffix in (".md", ".markdown", ".txt"):
        return path.read_text(encoding="utf-8")
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise ValueError(f"Unsupported input format: {suffix} ({path.name})")
    return extractor(path)
